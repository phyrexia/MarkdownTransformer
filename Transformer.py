import os
import sys
import logging
import uuid
import time
import argparse
import hashlib
from pptx import Presentation
from docx import Document

# Initialize logging
def setup_logging(debug=False):
    log_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'transformer.log')
    level = logging.DEBUG if debug else logging.INFO
    logging.basicConfig(
        filename=log_file,
        level=level,
        format='%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S',
        force=True
    )
    # Also log to console
    console = logging.StreamHandler()
    console.setLevel(level)
    logging.getLogger('').addHandler(console)

# SSL Setup for Mac (Avoid monkey-patching if possible)
def setup_ssl():
    """Attempts to fix SSL certificate issues on macOS using certifi bundle."""
    try:
        import certifi
        os.environ['SSL_CERT_FILE'] = certifi.where()
        os.environ['REQUESTS_CA_BUNDLE'] = certifi.where()
        logging.info("SSL certificates configured using certifi.")
    except ImportError:
        logging.warning("certifi not found. SSL issues might occur on macOS.")
    
    # Still keep the environment var for PyTorch
    os.environ["PYTORCH_ENABLE_MPS_FALLBACK"] = "1"

# Import marker-pdf components
try:
    from marker.converters.pdf import PdfConverter
    from marker.models import create_model_dict
    from marker.output import text_from_rendered
    MARKER_AVAILABLE = True
except ImportError as e:
    logging.error(f"Failed to import marker: {e}")
    MARKER_AVAILABLE = False
except Exception as e:
    logging.error(f"Unexpected error importing marker: {e}")
    MARKER_AVAILABLE = False

# Global variable to hold models
MARKER_MODELS = None

class DummyTableResult:
    """Dummy table result with cells attribute"""
    def __init__(self):
        self.cells = []
        self.rows = []
        self.cols = []
    def __getattr__(self, name):
        return []

class DummyTableRecModel:
    """Dummy table recognition model for MPS compatibility"""
    def __call__(self, images, *args, **kwargs):
        if isinstance(images, list):
            return [DummyTableResult() for _ in images]
        return [DummyTableResult()]
    def __getattr__(self, name):
        return lambda *args, **kwargs: []

import subprocess

class DocumentConverter:
    def __init__(self, device=None):
        self.models = None
        self.device = device
        self._init_device()

    def _init_device(self):
        if self.device:
            return
        try:
            import torch
            if torch.backends.mps.is_available():
                self.device = "mps"
            elif torch.cuda.is_available():
                self.device = "cuda"
            else:
                self.device = "cpu"
        except ImportError:
            self.device = "cpu"

    def _pptx_to_pdf_macos(self, pptx_path):
        """Converts PPTX to PDF using Microsoft PowerPoint via AppleScript on macOS."""
        abs_pptx = os.path.abspath(pptx_path)
        pdf_path = os.path.splitext(abs_pptx)[0] + "_temp.pdf"
        
        # AppleScript to save as PDF
        script = f'''
        tell application "Microsoft PowerPoint"
            open POSIX file "{abs_pptx}"
            set activePres to active presentation
            save activePres in POSIX file "{pdf_path}" as save as PDF
            close activePres saving no
        end tell
        '''
        try:
            logging.info(f"Converting PPTX to temp PDF using PowerPoint: {pptx_path}")
            subprocess.run(['osascript', '-e', script], check=True, capture_output=True)
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception as e:
            logging.error(f"PowerPoint PDF conversion failed: {e}")
        return None

    def _save_page_snapshots(self, filepath, image_dir, prefix="page", max_pages=None):
        """Renders each page of a PDF as optimized JPEG snapshots."""
        try:
            import fitz
            logging.info(f"Generating optimized snapshots for: {filepath}")
            doc = fitz.open(filepath)
            os.makedirs(image_dir, exist_ok=True)
            
            snapshots = []
            num_pages = len(doc)
            if max_pages:
                num_pages = min(num_pages, max_pages)
                
            for i in range(num_pages):
                page = doc.load_page(i)
                # 1.5x zoom (approx 150-200 DPI) is enough for thumbnails and saves >30% space
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                fname = f"{prefix}_{i+1}_snapshot.jpg"
                pix.save(os.path.join(image_dir, fname), "jpg", jpg_quality=80)
                snapshots.append(fname)
            doc.close()
            return snapshots
        except Exception as e:
            logging.error(f"Failed to generate snapshots: {e}")
            return []

    def _get_image_hash(self, data):
        """Returns SHA-256 hash of image data."""
        return hashlib.sha256(data).hexdigest()

    def _save_pptx_image(self, image_blob, ext, image_dir, image_dir_name, shape_name, slide_idx, shape_id, registry):
        """Check hash and save PPTX image if unique, otherwise return existing path."""
        img_hash = self._get_image_hash(image_blob)
        if img_hash in registry:
            return f"\n![{shape_name}]({registry[img_hash]})\n"
        
        fname = f"slide_{slide_idx+1}_{shape_id}_{uuid.uuid4().hex[:4]}.{ext}"
        fpath = os.path.join(image_dir, fname)
        with open(fpath, "wb") as f:
            f.write(image_blob)
        
        rel_path = f"{image_dir_name}/{fname}"
        registry[img_hash] = rel_path
        return f"\n![{shape_name}]({rel_path})\n"

    def _save_pdf_image(self, pil_image, image_dir, filename, registry):
        """Check hash and save PDF image if unique (using JPEG for consistency)."""
        import io
        img_byte_arr = io.BytesIO()
        if pil_image.mode in ("RGBA", "P"):
            pil_image = pil_image.convert("RGB")
        
        # We use a fixed quality and format for hashing consistency
        pil_image.save(img_byte_arr, format='JPEG', quality=85)
        img_hash = self._get_image_hash(img_byte_arr.getvalue())
        
        # Important: marker-pdf might use different extensions in its text.
        # We normalize everything to .jpg for deduplication storage.
        base_name = os.path.splitext(filename)[0]
        new_filename = f"{base_name}.jpg"
        
        if img_hash in registry:
            return registry[img_hash], False # Path, IsNew
        
        pil_image.save(os.path.join(image_dir, new_filename), "JPEG", quality=85)
        registry[img_hash] = new_filename
        return new_filename, True

    def load_models(self, force_cpu=False):
        """Load or reload models, optionally forcing CPU mode."""
        if not MARKER_AVAILABLE:
            return False
        
        device_to_use = "cpu" if force_cpu else self.device
        
        # If already loaded on correct device, return
        if self.models is not None and getattr(self, '_current_model_device', None) == device_to_use:
            return True
        
        logging.info(f"Loading AI models on {device_to_use}...")
        try:
            self.models = create_model_dict(device=device_to_use)
            self._current_model_device = device_to_use
            # Apply MPS workaround for tables
            if device_to_use == "mps" and "table_rec_model" in self.models:
                self.models["table_rec_model"] = DummyTableRecModel()
                logging.info("Using dummy table model for MPS compatibility.")
            return True
        except Exception as e:
            logging.error(f"Failed to load models on {device_to_use}: {e}")
            return False

    def _get_metadata_frontmatter(self, filename, metadata=None):
        """Generates YAML frontmatter with document metadata."""
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        fm = "---\n"
        fm += f"title: {filename}\n"
        fm += f"processed_at: {timestamp}\n"
        
        if metadata:
            for key, value in metadata.items():
                if value:
                    # Sanitize value to avoid breaking YAML
                    clean_value = str(value).replace('"', '\\"').strip()
                    fm += f"{key}: \"{clean_value}\"\n"
        
        fm += "---\n\n"
        return fm

    def convert_pdf(self, filepath, max_pages=None):
        """Convert PDF with prioritized snapshots, deduplication, and CPU fallback."""
        # Step 1: Setup Directories
        base_name = os.path.splitext(os.path.basename(filepath))[0]
        image_dir_name = f"{base_name}_images"
        image_dir = os.path.join(os.path.dirname(filepath), image_dir_name)
        os.makedirs(image_dir, exist_ok=True)
        
        # Step 1.5: Handle page subsetting if max_pages is set
        effective_path = filepath
        temp_subset_pdf = None
        doc_metadata = {}
        
        try:
            import fitz
            # Extract metadata from the original file before any subsetting
            with fitz.open(filepath) as meta_doc:
                doc_metadata = {
                    "original_author": meta_doc.metadata.get('author'),
                    "original_title": meta_doc.metadata.get('title'),
                    "page_count": meta_doc.page_count
                }
            
            if max_pages:
                logging.info(f"Creating temporary subset PDF for first {max_pages} pages")
                src = fitz.open(filepath)
                if len(src) > max_pages:
                    temp_subset_pdf = os.path.join(os.path.dirname(filepath), f"temp_subset_{uuid.uuid4().hex}.pdf")
                    dst = fitz.open()
                    dst.insert_pdf(src, from_page=0, to_page=max_pages-1)
                    dst.save(temp_subset_pdf)
                    dst.close()
                    effective_path = temp_subset_pdf
                src.close()
        except Exception as e:
            logging.error(f"Failed to handle PDF setup: {e}")

        # Step 2: Page Snapshots (Highest priority for faithful copy)
        snapshots = self._save_page_snapshots(filepath, image_dir, max_pages=max_pages)
        image_registry = {} # hash -> filename
        
        # Step 3: Marker extraction
        text = ""
        success = False
        
        # Try primary device (GPU/MPS) first
        if MARKER_AVAILABLE and self.load_models():
            try:
                logging.info(f"Extracting text via marker (Primary: {self.device}): {effective_path}")
                converter = PdfConverter(artifact_dict=self.models)
                rendered = converter(effective_path)
                ext_text, _, images = text_from_rendered(rendered)
                text = ext_text
                # Step 3a: Deduplicate and save images
                if images:
                    for filename, image in images.items():
                        saved_name, is_new = self._save_pdf_image(image, image_dir, filename, image_registry)
                        # Always replace to ensure extension matches (normalized to .jpg)
                        text = text.replace(filename, saved_name)
                success = True
            except Exception as e:
                logging.warning(f"Primary extraction failed ({e}). Retrying on CPU...")
                
        # CPU Fallback if primary failed or not successful
        if not success and MARKER_AVAILABLE and self.load_models(force_cpu=True):
            try:
                logging.info(f"Extracting text via marker (Fallback: CPU): {effective_path}")
                converter = PdfConverter(artifact_dict=self.models)
                rendered = converter(effective_path)
                ext_text, _, images = text_from_rendered(rendered)
                text = ext_text
                if images:
                    for filename, image in images.items():
                        saved_name, is_new = self._save_pdf_image(image, image_dir, filename, image_registry)
                        text = text.replace(filename, saved_name)
                success = True
            except Exception as e:
                logging.error(f"CPU Fallback extraction failed: {e}")
                text = f"\n> [!CAUTION]\n> AI extraction failed on both GPU and CPU: {e}\n\n"
        
        if not success and not text:
            text = "\n> [!NOTE]\n> AI extraction skipped (models not available).\n\n"

        # Step 4: Optimized Snapshot Interleaving (Single Pass)
        import re
        
        # Create a mapping for quick lookup: 1 -> "image_filename.jpg"
        # snapshots list is 0-indexed, so page 1 is snapshots[0]
        page_map = {i+1: snap for i, snap in enumerate(snapshots)}
        inserted_pages = set()

        def replacement_handler(match):
            try:
                p_num = int(match.group(1))
                if p_num in page_map:
                    inserted_pages.add(p_num)
                    img_md = f'\n![Page {p_num} Preview]({image_dir_name}/{page_map[p_num]})\n'
                    return f'{img_md}{match.group(0)}'
            except:
                pass
            return match.group(0)

        # Regex to find <span id="page-X-0"></span>
        pattern = re.compile(r'<span id="page-(\d+)-0"></span>')
        text = pattern.sub(replacement_handler, text)
        
        # Handle non-anchored pages (prepend them)
        sorted_pages = sorted(page_map.keys())
        prepend_text = ""
        for p_num in sorted_pages:
            if p_num not in inserted_pages:
                img_md = f'\n![Page {p_num} Preview]({image_dir_name}/{page_map[p_num]})\n'
                prepend_text += f"## Page {p_num} (Snapshot)\n{img_md}\n"
        
        if prepend_text:
            text = prepend_text + text

        # Step 5: Add Frontmatter
        frontmatter = self._get_metadata_frontmatter(base_name, doc_metadata)
        text = frontmatter + text

        # Step 6: Cleanup temporary files
        if temp_subset_pdf and os.path.exists(temp_subset_pdf):
            try:
                os.remove(temp_subset_pdf)
            except:
                pass

        return text

    def convert_pptx(self, filepath):
        """Convert PPTX to Markdown with interleaved slide snapshots and deduplicated images."""
        try:
            base_name = os.path.splitext(os.path.basename(filepath))[0]
            image_dir_name = f"{base_name}_images"
            image_dir = os.path.join(os.path.dirname(filepath), image_dir_name)
            os.makedirs(image_dir, exist_ok=True)
            
            # Extract Metadata
            prs = Presentation(filepath)
            doc_metadata = {
                "original_author": prs.core_properties.author,
                "original_title": prs.core_properties.title,
                "created": prs.core_properties.created,
                "slide_count": len(prs.slides)
            }

            # Step 1: Slide Snapshots (via temp PDF on macOS)
            pdf_temp = self._pptx_to_pdf_macos(filepath)
            snapshots = []
            if pdf_temp:
                snapshots = self._save_page_snapshots(pdf_temp, image_dir, prefix="slide")
                if os.path.exists(pdf_temp):
                    os.remove(pdf_temp)

            # Step 2: Normal extraction via python-pptx
            md_content = ""
            image_registry = {} # hash -> rel_path
            
            def get_images_from_shape(slide_idx, shape, slide_image_dir, registry):
                """Recursive function to find images within shapes, groups, and placeholders."""
                imgs = []
                if shape.shape_type == 6: # GROUP
                    for s in shape.shapes: imgs.extend(get_images_from_shape(slide_idx, s, slide_image_dir, registry))
                elif shape.shape_type == 13: # PICTURE
                    try:
                        img = shape.image
                        if img.ext.lower() in ['wmf', 'emz']: return []
                        img_md = self._save_pptx_image(img.blob, img.ext.lower(), image_dir, image_dir_name, shape.name, slide_idx, shape.shape_id, registry)
                        imgs.append(img_md)
                    except Exception as e: logging.warning(f"Image error: {e}")
                elif shape.is_placeholder:
                    try: 
                        if hasattr(shape, 'image'): imgs.extend(get_images_from_shape(slide_idx, shape, slide_image_dir, registry))
                    except: pass
                return imgs

            for i, slide in enumerate(prs.slides):
                md_content += f"\n## Slide {i+1}\n\n"
                
                # Insert Snapshot Interleaved
                if i < len(snapshots):
                    md_content += f"![Slide {i+1} Preview]({image_dir_name}/{snapshots[i]})\n\n"
                
                # Title
                if slide.shapes.title:
                    md_content += f"# {slide.shapes.title.text.strip()}\n\n"
                
                # Body Content
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                md_content += f"{'  '*para.level}- {para.text.strip()}\n"
                    for img_md in get_images_from_shape(i, shape, image_dir, image_registry):
                        md_content += img_md
                
                md_content += "\n---\n"
            
            # Step 3: Prepend Frontmatter
            frontmatter = self._get_metadata_frontmatter(base_name, doc_metadata)
            md_content = frontmatter + md_content

            return md_content
        except Exception as e:
            logging.error(f"PPTX Error: {e}")
            return f"Error during PPTX conversion: {e}"

    def convert_docx(self, filepath):
        """Convert DOCX to Markdown."""
        try:
            doc = Document(filepath)
            
            doc_metadata = {
                "original_author": doc.core_properties.author,
                "original_title": doc.core_properties.title,
                "created": doc.core_properties.created,
            }
            
            md_content = ""
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                style_name = para.style.name.lower()
                if 'heading 1' in style_name:
                    md_content += f"# {text}\n\n"
                elif 'heading 2' in style_name:
                    md_content += f"## {text}\n\n"
                elif 'heading 3' in style_name:
                    md_content += f"### {text}\n\n"
                elif 'list' in style_name:
                    md_content += f"- {text}\n"
                else:
                    md_content += f"{text}\n\n"
            
            frontmatter = self._get_metadata_frontmatter(os.path.basename(filepath), doc_metadata)
            return frontmatter + md_content
        except Exception as e:
            logging.error(f"DOCX Error: {e}")
            return f"Error during DOCX conversion: {e}"

def main():
    parser = argparse.ArgumentParser(description="Convert Document files (PDF, PPTX, DOCX) to Markdown.")
    parser.add_argument("files", nargs="+", help="File paths to convert.")
    parser.add_argument("--device", choices=["cpu", "mps", "cuda"], help="Force specific device.")
    parser.add_argument("--debug", action="store_true", help="Enable debug logging.")
    parser.add_argument("--pages", type=int, help="Limit number of pages to process.")
    args = parser.parse_args()

    setup_logging(args.debug)
    setup_ssl()
    
    converter = DocumentConverter(device=args.device)
    total = len(args.files)
    
    for idx, filepath in enumerate(args.files, 1):
        if not os.path.exists(filepath):
            logging.error(f"File not found: {filepath}")
            continue
            
        file_basename = os.path.basename(filepath)
        logging.info(f"[{idx}/{total}] Processing: {file_basename}")
        print(f"[{idx}/{total}] Converting {file_basename}...")
        
        start_time = time.time()
        ext = os.path.splitext(filepath)[1].lower()
        md_output = None
        
        if ext == ".pdf":
            md_output = converter.convert_pdf(filepath, max_pages=args.pages)
        elif ext == ".pptx":
            md_output = converter.convert_pptx(filepath)
        elif ext == ".docx":
            md_output = converter.convert_docx(filepath)
        else:
            logging.warning(f"Unsupported extension {ext}")
            continue

        if md_output:
            out_path = os.path.splitext(filepath)[0] + ".md"
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(md_output)
            duration = time.time() - start_time
            logging.info(f"SUCCESS: {file_basename} in {duration:.2f}s")
            print(f"✓ Created {os.path.basename(out_path)}")
        else:
            logging.error(f"FAILURE: {file_basename}")

if __name__ == "__main__":
    main()
